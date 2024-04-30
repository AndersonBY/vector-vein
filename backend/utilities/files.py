# -*- coding: utf-8 -*-
# @Author: Bi Ying
# @Date:   2023-05-15 11:17:29
# @Last Modified by:   Bi Ying
# @Last Modified time: 2024-04-29 20:18:56
import json
import shutil
import zipfile
from pathlib import Path

import pypdf
import mammoth
import openpyxl
import pathspec
import pandas as pd
from pptx import Presentation

from utilities.print_utils import mprint_error


CODEC_TEST_LIST = [
    "utf-8-sig",
    "utf-8",
    "gbk",
    "gb2312",
    "utf-16",
    "utf-16-le",
    "utf-16-be",
    "utf-32",
    "utf-32-le",
    "utf-32-be",
]


def try_load_json_file(file_path: str, default=dict):
    if Path(file_path).exists():
        try:
            with open(file_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except Exception as e:
            mprint_error(f"Failed to load json file: {file_path}, {e}")
    return default()


def rename_path(path: str | Path):
    path = Path(path)

    for item in path.iterdir():
        if item.is_file() or item.is_dir():
            try:
                new_name = item.name.encode("cp437").decode("gbk")
            except Exception:
                new_name = item.name

            new_path = item.with_name(new_name)
            item.rename(new_path)

            if new_path.is_dir():
                rename_path(new_path)


def read_gitignore(folder_path, extra_ignore_patterns=None):
    gitignore_path = folder_path / ".gitignore"
    ignore_lines = []

    if gitignore_path.exists():
        with gitignore_path.open("r", encoding="utf-8") as file:
            ignore_lines.extend(line.strip() for line in file)

    if extra_ignore_patterns:
        ignore_lines.extend(extra_ignore_patterns)

    spec = pathspec.PathSpec.from_lines("gitwildmatch", ignore_lines)
    return spec


def is_ignored(path, spec, folder_path):
    relative_path = path.relative_to(folder_path)
    # Use pathspec to check if the file should be ignored
    return spec.match_file(relative_path.as_posix())


def read_folder(folder_path, ignore_list):
    contents = []
    for path in sorted(folder_path.rglob("*")):
        if path.is_file() and not is_ignored(path, ignore_list, folder_path):
            try:
                file_content = read_file_content(filename=path.name, local_file=path, read_zip=False)
                rel_path = path.relative_to(folder_path)
                contents.append(f"# {rel_path}\n```\n{file_content}\n```\n")
            except Exception as e:
                print(f"Could not read file {path}: {e}")
    return contents


def read_zip_contents(zip_file_path: str | Path, extract_path: str | Path = "/tmp") -> Path:
    if not zipfile.is_zipfile(zip_file_path):
        raise FileNotFoundError(f"The file at {zip_file_path} is not a valid zip file.")

    zip_path = Path(zip_file_path)
    extract_to_path = Path(extract_path) / zip_path.stem

    extract_to_path.mkdir(parents=True, exist_ok=True)

    with zipfile.ZipFile(zip_file_path, "r") as zip_ref:
        zip_ref.extractall(extract_to_path)

    rename_path(extract_to_path)

    folder_path = Path(extract_to_path)
    extra_ignore_list = [".git", ".gitignore", ".gitattributes"]
    ignore_list = read_gitignore(folder_path, extra_ignore_list)

    contents = read_folder(folder_path, ignore_list)
    return "\n".join(contents)


def read_file_content(local_file: str | Path, read_zip: bool = False):
    filename = Path(local_file).name
    if filename.endswith(".docx"):
        with open(local_file, "rb") as docx_file:
            docx_data = mammoth.convert_to_markdown(docx_file)
            markdown_text = docx_data.value
            return markdown_text
    elif filename.endswith(".pdf"):
        with open(local_file, "rb") as pdf_file_obj:
            pdf_reader = pypdf.PdfReader(pdf_file_obj)
            pdf_contents = [page.extract_text() for page in pdf_reader.pages]
            return "\n\n".join(pdf_contents)
    elif filename.endswith(".pptx"):
        ppt = Presentation(local_file)
        ppt_contents = []
        for slide in ppt.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    ppt_contents.append(shape.text.strip())
        return "\n\n".join(ppt_contents)
    elif filename.endswith(".xlsx"):
        try:
            df = pd.read_excel(local_file, engine="openpyxl")
            return df.to_csv(index=False)
        except Exception as e:
            mprint_error(e)
            wb = openpyxl.load_workbook(local_file, data_only=True)
            ws = wb.active
            csv_contents = []
            for row in ws.rows:
                csv_contents.append(",".join([str(cell.value) if cell.value else "" for cell in row]))
            return "\n\n".join(csv_contents)
    elif filename.endswith((".txt", ".md", ".html", ".json", ".csv", ".srt")):
        for codec in CODEC_TEST_LIST:
            try:
                with open(local_file, "r", encoding=codec) as txt_file:
                    txt_contents = txt_file.read()
                    return txt_contents
            except Exception as e:
                mprint_error(e)
        else:
            mprint_error("Failed to decode file")
            return ""
    elif filename.endswith(".zip") and read_zip:
        content = read_zip_contents(local_file)
        return content
    else:
        for codec in CODEC_TEST_LIST:
            try:
                with open(local_file, "r", encoding=codec) as txt_file:
                    txt_contents = txt_file.read()
                    return txt_contents
            except Exception as e:
                mprint_error(e)
        else:
            mprint_error("Failed to decode file")
            return txt_contents


def get_files_contents(files: list):
    results = []
    for file in files:
        results.append(read_file_content(file, read_zip=True))

    return results


def copy_file(src, dst):
    if Path(src).exists():
        Path(dst).parent.mkdir(parents=True, exist_ok=True)
        shutil.copy(src, dst)
        return True
    return False
